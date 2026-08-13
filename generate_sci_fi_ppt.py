import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_scifi_theme(slide, is_title=False):
    # Set dark ethereal background (Deep Space Blue/Purple)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(12, 18, 40)  # Dark ethereal blue-black
    
    # We will style shapes individually since pptx doesn't allow easy master slide overriding from scratch without a template

def format_title(shape, text, font_size=44):
    shape.text = text
    if not shape.text_frame.paragraphs:
        return
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Century Gothic'
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0, 230, 255)  # Glowing Cyan

def format_body(shape, text, font_size=20, is_bullet=True):
    tf = shape.text_frame
    tf.text = ""  # clear
    
    for line in text.split('\n'):
        p = tf.add_paragraph()
        p.text = line.strip()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(220, 220, 240)  # Soft white
        p.level = 0 if is_bullet else None
        
def create_ppt():
    prs = Presentation()
    
    # SLIDE 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    apply_scifi_theme(slide, is_title=True)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    format_title(title, "ASCENSION: EARLY BRAIN TUMOR DETECTION", 48)
    
    subtitle.text = "AI Meets the Ethereal | Bridging Healthcare and the Future"
    for p in subtitle.text_frame.paragraphs:
        p.font.name = 'Century Gothic'
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 215, 0) # Ethereal Gold
        
    # SLIDE 2: Vision
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    apply_scifi_theme(slide)
    
    format_title(slide.shapes.title, "The Vision")
    body_text = (
        "An end-to-end, deep learning-powered medical imaging web application.\n"
        "Classifies brain MRI scans into four distinct categories:\n"
        "  - Glioma\n"
        "  - Meningioma\n"
        "  - Pituitary Tumor\n"
        "  - No Tumor\n"
        "Provides seamless user experience with real-time AI inference."
    )
    format_body(slide.placeholders[1], body_text)
    
    # SLIDE 3: Data Cosmos
    slide = prs.slides.add_slide(bullet_slide_layout)
    apply_scifi_theme(slide)
    format_title(slide.shapes.title, "The Data Cosmos")
    body_text = (
        "Trained on a meticulously curated Mega-Dataset of 13,100+ MRI scans.\n"
        "Aggressive real-time data augmentations applied:\n"
        "  - ±20° Rotations\n"
        "  - Spatial shifting and zooming\n"
        "  - Horizontal flipping\n"
        "Guarantees that the AI learns fundamental structures, not just artifacts."
    )
    format_body(slide.placeholders[1], body_text)
    
    # SLIDE 4: Neural Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    apply_scifi_theme(slide)
    format_title(slide.shapes.title, "The Neural Architectures")
    body_text = (
        "DenseNet121 (Fine-Tuned): Deepest 40 layers unfrozen to adapt to MRI textural gradients.\n"
        "Explainable AI (Grad-CAM): Custom Pure-NumPy gradient mapper reveals exactly what the AI sees.\n"
        "Clinical Diversity: Merged multiple datasets to eliminate overfitting."
    )
    format_body(slide.placeholders[1], body_text)

    # SLIDE 5: Zenith of Accuracy
    slide = prs.slides.add_slide(bullet_slide_layout)
    apply_scifi_theme(slide)
    format_title(slide.shapes.title, "The Zenith of Accuracy")
    body_text = (
        "95% Overall Accuracy achieved after deep fine-tuning.\n"
        "98% Recall for 'No Tumor' (Prioritizing safe diagnosis).\n"
        "100% Recall for Pituitary Tumors (Extreme detection confidence).\n"
        "A highly robust, production-ready AI."
    )
    format_body(slide.placeholders[1], body_text)

    # SLIDE 6: Ethereal Engine (Tech Stack)
    slide = prs.slides.add_slide(bullet_slide_layout)
    apply_scifi_theme(slide)
    format_title(slide.shapes.title, "The Ethereal Engine (Tech Stack)")
    body_text = (
        "Frontend: React + Vite (Blazing fast, Glassmorphism UI, interactive charts).\n"
        "Backend API: Python FastAPI (Asynchronous processing).\n"
        "Inference Engine: ONNX Runtime (Replaces massive 1GB TensorFlow payloads with <50MB hyper-optimized graphs).\n"
        "Deployments: Vercel (Frontend) & Render (Backend)."
    )
    format_body(slide.placeholders[1], body_text)
    
    # SLIDE 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_scifi_theme(slide, is_title=True)
    format_title(slide.shapes.title, "The Future is Here", 54)
    subtitle = slide.placeholders[1]
    subtitle.text = "Thank You\nDesigned by Ashu"
    for p in subtitle.text_frame.paragraphs:
        p.font.name = 'Century Gothic'
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0, 230, 255)

    prs.save('Sci_Fi_Brain_Tumor_Presentation.pptx')
    print("Presentation saved as Sci_Fi_Brain_Tumor_Presentation.pptx")

if __name__ == '__main__':
    create_ppt()
